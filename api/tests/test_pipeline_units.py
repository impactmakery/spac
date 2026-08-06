import io

import pytest

# ---------- upload validation ----------


def _pdf_bytes() -> bytes:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _docx_bytes(text="שלום עולם. Budget guidance for 2026.") -> bytes:
    import docx

    doc = docx.Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_validate_upload_accepts_known_types():
    from app.services.uploads import validate_upload

    ext, ct = validate_upload("Doc.PDF", _pdf_bytes(), "application/pdf")
    assert (ext, ct) == ("pdf", "application/pdf")
    ext, _ = validate_upload("word.docx", _docx_bytes(), "application/octet-stream")
    assert ext == "docx"


def test_validate_upload_rejects_bad_extension_and_magic():
    from fastapi import HTTPException

    from app.services.uploads import validate_upload

    with pytest.raises(HTTPException) as e:
        validate_upload("script.exe", b"MZ....", "application/octet-stream")
    assert e.value.status_code == 415

    # extension says pdf but bytes are not a pdf
    with pytest.raises(HTTPException) as e:
        validate_upload("fake.pdf", b"not a pdf at all", "application/pdf")
    assert e.value.status_code == 415


def test_validate_upload_rejects_oversize():
    from fastapi import HTTPException

    from app.services.uploads import MAX_UPLOAD_BYTES, validate_upload

    big = b"%PDF-" + b"0" * (MAX_UPLOAD_BYTES + 1)
    with pytest.raises(HTTPException) as e:
        validate_upload("big.pdf", big, "application/pdf")
    assert e.value.status_code == 413


# ---------- chunking ----------


def test_chunk_text_bounds_and_overlap():
    from app.rag.chunking import chunk_text, token_len

    text = " ".join(f"word{i}" for i in range(3000))
    chunks = chunk_text(text, max_tokens=800, overlap=150)
    assert len(chunks) > 1
    for c in chunks:
        assert token_len(c) <= 800
    # consecutive chunks share content (overlap)
    assert chunks[0][-40:] != chunks[1][-40:]
    assert any(w in chunks[1] for w in chunks[0].split()[-30:])


def test_chunk_text_short_input_single_chunk():
    from app.rag.chunking import chunk_text

    assert chunk_text("hello world") == ["hello world"]
    assert chunk_text("   ") == []


# ---------- embeddings ----------


def test_fake_embeddings_deterministic_unit_norm():
    import math

    from app.rag.embeddings import get_embedding_provider

    provider = get_embedding_provider()
    assert type(provider).__name__ == "FakeEmbeddings"  # no OPENAI_API_KEY in tests
    [a1], [a2] = provider.embed(["שלום"]), provider.embed(["שלום"])
    [b] = provider.embed(["something else"])
    assert a1 == a2 and a1 != b
    assert len(a1) == 1536
    assert abs(math.sqrt(sum(x * x for x in a1)) - 1.0) < 1e-6


# ---------- extraction ----------


def test_extract_docx_pptx_xlsx_pdf():
    from app.rag.extract import extract_text

    assert "Budget guidance" in extract_text(_docx_bytes(), "docx")

    import pptx

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly plan"
    buf = io.BytesIO()
    prs.save(buf)
    assert "Quarterly plan" in extract_text(buf.getvalue(), "pptx")

    import openpyxl

    wb = openpyxl.Workbook()
    wb.active.append(["Item", "Cost"])
    wb.active.append(["Chairs", 1200])
    buf = io.BytesIO()
    wb.save(buf)
    text = extract_text(buf.getvalue(), "xlsx")
    assert "Chairs" in text and "1200" in text

    # blank pdf extracts to empty-ish, but must not raise
    assert isinstance(extract_text(_pdf_bytes(), "pdf"), str)


def test_extract_corrupt_raises():
    from app.rag.extract import ExtractionError, extract_text

    with pytest.raises(ExtractionError):
        extract_text(b"garbage bytes", "docx")


# ---------- structure-aware chunking ----------


def test_chunks_break_at_paragraph_boundaries():
    """A chunk that ends mid-sentence embeds poorly and reads badly when cited."""
    from app.rag.chunking import chunk_text, token_len

    paragraphs = [f"Paragraph {i}. " + " ".join(f"word{j}" for j in range(120))
                  for i in range(12)]
    chunks = chunk_text("\n\n".join(paragraphs), max_tokens=400, overlap=60)

    assert len(chunks) > 1
    for c in chunks:
        assert token_len(c) <= 400
        # every chunk starts at a paragraph start, never mid-sentence
        assert c.startswith("Paragraph")


def test_title_is_prepended_to_every_chunk():
    from app.rag.chunking import chunk_text

    body = "\n\n".join(f"Section {i} body text here." for i in range(5))
    chunks = chunk_text(body, max_tokens=20, overlap=5, title="Waste Guidelines")
    assert len(chunks) > 1
    assert all(c.startswith("Waste Guidelines") for c in chunks)
    # the title costs budget but must never displace content
    joined = " ".join(chunks)
    assert all(f"Section {i}" in joined for i in range(5))


def test_overlap_carries_context_across_the_seam():
    from app.rag.chunking import chunk_text

    paragraphs = [f"Block{i} " + " ".join(["filler"] * 40) for i in range(8)]
    chunks = chunk_text("\n\n".join(paragraphs), max_tokens=200, overlap=80)
    assert len(chunks) > 1
    # the start of a later chunk repeats the end of the previous one
    first_blocks = {line.split()[0] for line in chunks[0].split("\n\n")}
    second_blocks = {line.split()[0] for line in chunks[1].split("\n\n")}
    assert first_blocks & second_blocks


def test_single_oversized_block_is_still_split():
    from app.rag.chunking import chunk_text, token_len

    wall = " ".join(f"word{i}" for i in range(2000))  # no paragraph breaks at all
    chunks = chunk_text(wall, max_tokens=300, overlap=50)
    assert len(chunks) > 1
    assert all(token_len(c) <= 300 for c in chunks)


def test_short_document_stays_one_chunk():
    from app.rag.chunking import chunk_text

    assert chunk_text("A short procedure note.") == ["A short procedure note."]
    assert chunk_text("   ") == []


# ---------- OCR for scanned PDFs ----------


def _image_pdf(text="Scanned municipal notice 2026", size=(1000, 260)) -> bytes:
    """A PDF whose page is an image — exactly what a scanner produces."""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(img)
    draw.text((20, 100), text, fill="black")
    buf = io.BytesIO()
    img.save(buf, format="PDF")
    return buf.getvalue()


def test_scanned_pdf_has_no_text_layer():
    """The premise: pypdf alone returns nothing for a scan, which is why these
    files were being marked 'not indexable'."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(_image_pdf()))
    assert "".join(p.extract_text() or "" for p in reader.pages).strip() == ""


def test_ocr_never_raises_when_tesseract_is_missing():
    """OCR is a best-effort improvement on a document that already extracted to
    nothing. A missing binary must not turn an empty document into a failed job."""
    from app.rag.extract import _ocr_pdf, extract_text

    assert isinstance(_ocr_pdf(_image_pdf()), str)
    assert isinstance(extract_text(_image_pdf(), "pdf"), str)


def test_ocr_is_not_attempted_on_a_pdf_that_has_real_text(monkeypatch):
    from app.rag import extract as extract_mod

    called = False

    def spy(content):
        nonlocal called
        called = True
        return ""

    monkeypatch.setattr(extract_mod, "_ocr_pdf", spy)
    text = extract_mod.extract_text(_text_pdf(), "pdf")
    assert "Budget planning cycle" in text
    assert not called, "a born-digital PDF must not pay the OCR cost"


def test_failed_ocr_does_not_destroy_a_weak_text_layer(monkeypatch):
    """A mixed document — a real report with scanned appendices — must keep its
    text layer even when the OCR pass returns nothing."""
    from app.rag import extract as extract_mod

    monkeypatch.setattr(extract_mod, "MIN_CHARS_PER_PAGE", 10_000)  # force the attempt
    monkeypatch.setattr(extract_mod, "_ocr_pdf", lambda content: "")
    text = extract_mod.extract_text(_text_pdf(), "pdf")
    assert "Budget planning cycle" in text


def _text_pdf(line="Budget planning cycle begins in November for every municipality.") -> bytes:
    """A born-digital PDF with a genuine text layer.

    Assembled by hand rather than pulling in a PDF-authoring dependency for one
    fixture; the xref offsets are computed so pypdf parses it normally.
    """
    stream = f"BT /F1 12 Tf 72 720 Td ({line}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode() + body + b"\nendobj\n"

    xref_at = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_at}\n".encode()
    )
    out += b"%%EOF\n"
    return bytes(out)
