"""Text extraction per format. Raises ExtractionError on unreadable files.

Municipal archives are full of scanned paper: a PDF whose pages are images with
no text layer. Those extract to nothing and were previously marked
'not indexable'. When a PDF yields too little text to be plausible, the pages are
rasterised and run through Tesseract in Hebrew and English.
"""

import io
import logging

log = logging.getLogger(__name__)

# A born-digital page carries far more than this; anything less means the text
# layer is missing or decorative, and the real content is in the pixels.
MIN_CHARS_PER_PAGE = 40
OCR_MAX_PAGES = 40  # a 300-page scan would hold the worker for many minutes
OCR_DPI = 200  # below ~150 Hebrew diacritics and small print start to fail

# Ceiling on the bitmap a single page may be rendered to, in pixels.
#
# DPI alone bounds nothing: it is a multiplier on the page's own size, so an A4
# page at 200 dpi is a manageable 11 MB while an A0 plan is 182 MB — and
# Tesseract then takes its own copy. Municipal archives are full of scanned
# plans and posters, and the worker has a gigabyte. Eight megapixels is far
# more than Tesseract needs for readable text; anything larger is scaled down
# to fit rather than rendered at full size and hoped for.
OCR_MAX_PIXELS = 8_000_000
def _languages() -> str:
    """Which scripts Tesseract should try, from configuration.

    Hebrew and English were the original pair. Arabic was absent, so every
    scanned or picture-based document belonging to the two Arabic-speaking
    municipalities read as nothing at all — OCR working perfectly and
    recognising no words, which looks identical to a document with none.
    """
    from app.core.config import get_settings

    return get_settings().ocr_languages or "heb+eng"

# A Word file or slide deck whose text amounts to less than this is very likely
# a wrapper around pictures — a deck built in a design tool and pasted in, or a
# photographed page dropped into a document. Below the threshold the embedded
# images are read as well; above it they are decoration and not worth the time.
MIN_CHARS_OFFICE = 200

# Bounds on that pass. Municipal decks carry logos, headers and separators on
# every slide, and OCR'ing forty copies of a crest costs minutes and returns
# nothing. Small images are skipped, and only so many are read per document.
OCR_MIN_IMAGE_BYTES = 40_000
OCR_MAX_EMBEDDED_IMAGES = 20

# Where each Office format keeps its pictures inside the zip.
OFFICE_MEDIA_PREFIXES = ("word/media/", "ppt/media/", "xl/media/")


class ExtractionError(Exception):
    pass


IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "webp", "gif", "bmp", "tiff", "tif"}
TEXT_EXTENSIONS = {"txt", "csv", "md", "markdown", "json", "yaml", "yml"}


def extract_text(content: bytes, ext: str) -> str:
    ext = (ext or "").lower()
    try:
        if ext == "pdf":
            return _pdf(content)
        if ext == "docx":
            return _docx(content)
        if ext == "pptx":
            return _pptx(content)
        if ext == "xlsx":
            return _xlsx(content)
        if ext in IMAGE_EXTENSIONS:
            return _image(content)
        if ext in TEXT_EXTENSIONS:
            return _plain(content)
    except ExtractionError:
        raise
    except Exception as e:
        raise ExtractionError(str(e)) from e
    raise ExtractionError(f"unsupported extension: {ext}")


def _pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages = [(page.extract_text() or "") for page in reader.pages]
    text = "\n\n".join(pages)

    if pages and len(text.strip()) < MIN_CHARS_PER_PAGE * len(pages):
        scanned = _ocr_pdf(content)
        # Keep whichever is richer: a mixed document (a born-digital report with
        # scanned appendices) should not lose its real text layer to a failed
        # OCR pass, and a pure scan should not keep its empty one.
        if len(scanned.strip()) > len(text.strip()):
            return scanned
    return text


def _office_images(content: bytes) -> str:
    """Read the words in the pictures an Office file wraps around.

    Text extraction only sees text runs, so a deck whose slides are exported
    images contributes its title and nothing else — four documents in the
    first real corpus were exactly that, one of them a municipality's main
    presentation of its welfare services.

    Never raises: this is an improvement on a document that already extracted
    to almost nothing, and a failure here must not cost it the little it had.
    """
    import zipfile

    try:
        archive = zipfile.ZipFile(io.BytesIO(content))
    except Exception as e:  # noqa: BLE001
        log.warning("could not open Office file to read its images: %s", e)
        return ""

    parts: list[str] = []
    read = 0
    for info in sorted(archive.infolist(), key=lambda i: i.filename):
        if read >= OCR_MAX_EMBEDDED_IMAGES:
            log.info("stopped after %d embedded images", OCR_MAX_EMBEDDED_IMAGES)
            break
        name = info.filename
        if not name.startswith(OFFICE_MEDIA_PREFIXES):
            continue
        if name.rsplit(".", 1)[-1].lower() not in IMAGE_EXTENSIONS:
            continue  # embedded video, audio, or an unreadable blob
        if info.file_size < OCR_MIN_IMAGE_BYTES:
            continue  # a logo or an icon, not a page of text
        try:
            text = _image(archive.read(name))
        except Exception as e:  # noqa: BLE001 — one bad image is not a failure
            log.warning("OCR failed on embedded image %s: %s", name, e)
            continue
        read += 1
        if text.strip():
            parts.append(text.strip())
    return "\n\n".join(parts)


def _with_embedded_images(content: bytes, text: str) -> str:
    """Add what the pictures say when the document itself said almost nothing."""
    if len(text.strip()) >= MIN_CHARS_OFFICE:
        return text
    found = _office_images(content)
    if not found.strip():
        return text
    return f"{text}\n\n{found}".strip()


def _image(content: bytes) -> str:
    """Read the words in a picture.

    A screenshot of a circular or a photographed notice is a real way people
    share information, and without this the assistant sees only the filename.
    Returns empty rather than raising when OCR is unavailable — an unreadable
    image is still a valid attachment.
    """
    try:
        import io as _io

        import pytesseract
        from PIL import Image
    except ImportError:
        log.info("image OCR skipped: pytesseract/Pillow not installed")
        return ""
    try:
        with Image.open(_io.BytesIO(content)) as image:
            return pytesseract.image_to_string(image, lang=_languages())
    except Exception as e:  # noqa: BLE001
        log.warning("image OCR failed: %s", e)
        return ""


def _plain(content: bytes) -> str:
    """Plain text, CSV, Markdown. Tries UTF-8 first, then Hebrew's legacy encoding."""
    for encoding in ("utf-8-sig", "utf-8", "cp1255", "iso-8859-8"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def ocr_available() -> bool:
    """Whether both halves of the OCR path are installed: the Python bindings
    and the Tesseract binary itself."""
    try:
        import pypdfium2  # noqa: F401
        import pytesseract

        pytesseract.get_tesseract_version()
    except Exception:  # noqa: BLE001 — any failure means the path is unusable
        return False
    return True


def _render_scale(page: object) -> float:
    """Render scale for a page: OCR_DPI, or less if that would be enormous."""
    scale = OCR_DPI / 72
    try:
        width, height = page.get_size()  # type: ignore[attr-defined]
    except Exception:  # noqa: BLE001 — an unreadable size just uses the default
        return scale
    pixels = (width * scale) * (height * scale)
    if pixels > OCR_MAX_PIXELS:
        scale *= (OCR_MAX_PIXELS / pixels) ** 0.5
        log.info("OCR scaled a large page down to fit %d pixels", OCR_MAX_PIXELS)
    return scale


def _ocr_pdf(content: bytes) -> str:
    """Rasterise each page and read it with Tesseract.

    Deliberately never raises: OCR is a best-effort improvement on a document
    that already extracted to nothing, so a missing binary or a corrupt page
    must not turn a merely-empty document into a failed ingestion job.
    """
    try:
        import pypdfium2 as pdfium
        import pytesseract
    except ImportError:
        log.info("OCR skipped: pypdfium2/pytesseract not installed")
        return ""

    try:
        pdf = pdfium.PdfDocument(io.BytesIO(content))
    except Exception as e:  # noqa: BLE001
        log.warning("OCR skipped: could not open PDF for rasterising: %s", e)
        return ""

    out: list[str] = []
    try:
        count = min(len(pdf), OCR_MAX_PAGES)
        if len(pdf) > OCR_MAX_PAGES:
            log.warning(
                "OCR limited to the first %d of %d pages", OCR_MAX_PAGES, len(pdf)
            )
        for index in range(count):
            page = None
            image = None
            try:
                page = pdf[index]
                image = page.render(scale=_render_scale(page)).to_pil()
                out.append(pytesseract.image_to_string(image, lang=_languages()))
            except Exception as e:  # noqa: BLE001 — one bad page is not a failure
                log.warning("OCR failed on page %d: %s", index + 1, e)
            finally:
                # Release each page's bitmap before rendering the next one.
                # Without this the native buffers accumulate for the whole
                # document and a long scan takes the worker out of memory.
                if image is not None:
                    image.close()
                if page is not None:
                    page.close()
    finally:
        pdf.close()

    return "\n\n".join(part for part in out if part.strip())


def _docx(content: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(content))
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    text = "\n".join(part for part in parts if part.strip())
    return _with_embedded_images(content, text)


def _pptx(content: bytes) -> str:
    import pptx

    prs = pptx.Presentation(io.BytesIO(content))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    text = "\n".join(part for part in parts if part.strip())
    return _with_embedded_images(content, text)


def _xlsx(content: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(str(ws.title))
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    return _with_embedded_images(content, text)
